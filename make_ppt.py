from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 1. 初始化 PPT 引擎
prs = Presentation()


# 定义全局样式函数
def set_title_style(title_shape):
    title_shape.text_frame.paragraphs[0].font.name = '微软雅黑'
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # 深蓝色


def add_content_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]  # 使用“标题和内容”版式
    slide = prs.slides.add_slide(slide_layout)

    # 设置标题
    title = slide.shapes.title
    title.text = title_text
    set_title_style(title)

    # 设置正文内容
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullet_points[0]

    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = '微软雅黑'
        p.font.size = Pt(20)


# ==========================================
# 🚀 第一页：封面 (P1)
# ==========================================
slide_layout = prs.slide_layouts[0]  # 封面版式
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "数字孪生·古今共振\n长春市宽城区核心街区更新规划设计"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "答辩人：陈礼冲 (城乡规划 21级)\n指导教师：崔诚哲\n核心技术：数字孪生 / CV测度 / AIGC / 大模型计算"

# ==========================================
# 🚀 第二页：研究背景与痛点 (P2)
# ==========================================
add_content_slide(prs, "1. 研究背景与核心痛点", [
    "📍 规划区位：长春市宽城区伪满皇宫周边 150 公顷复合街区。",
    "⚠️ 核心痛点：",
    "  - 产业空心化：中车厂区产能外迁，空置率高达 40%。",
    "  - 社会老龄化：老旧小区密集，老龄化率达 30%，缺乏活力。",
    "  - 空间割裂化：厂区与铁路线阻断路网，历史风貌遭违建侵蚀。",
    "💡 破局思路：引入“数字孪生”与“AI”技术，拒绝大拆大建，实施精准“微更新”。"
])

# ==========================================
# 🚀 第三页：技术架构 (P3)
# ==========================================
add_content_slide(prs, "2. 多模态微更新决策系统架构", [
    "传统规划调研难以处理高维数据，本课题自主研发了全栈 Web 决策平台：",
    "🔹 数据感知层：百度街景 + 交通潮汐 + POI + 社交媒体文本。",
    "🔹 诊断分析层：",
    "  - CV 引擎 (DeepLabV3+)：提取绿视率(GVI)与天空开阔度(SVF)。",
    "  - NLP 引擎 (Jieba + LLM)：挖掘公众情感极性与空间需求。",
    "🔹 推演决策层：基于 Stable Diffusion + ControlNet 进行风貌约束生成。"
])

# ==========================================
# 🚀 第四页：空间诊断 (P4)
# ==========================================
add_content_slide(prs, "3. 空间病理学诊断：多元数据耦合", [
    "基于 WebGL 的高并发数字沙盘渲染：",
    "🌿 空间品质测度：精准定位“视觉杂乱度极高”的违建盲区与“绿视率极低”的生态洼地。",
    "🚦 潮汐与路网耦合：建立 24H 动态模型，识别长春站周边的断头路与贯通靶点。",
    "😠 情感图谱映射：将公众社交文本转化为空间热力图，精准锚定“负面情绪”极化区域（如批发市场），将居民诉求量化为空间干预依据。"
])

# ==========================================
# 🚀 第五页：AIGC 介入 (P5)
# ==========================================
add_content_slide(prs, "4. AIGC 赋能：重点地块微更新推演", [
    "针对 5 个核心衰败靶点，执行‘古今共振’风貌推演：",
    "🎯 物理空间约束：使用 ControlNet (Canny 算子) 提取现状透视骨架，拒绝 AI 随机幻觉。",
    "🎨 多元策略重绘：",
    "  - 策略 1：历史风貌修缮（伪满时期遗存的立面还原）。",
    "  - 策略 2：工业遗迹复兴（中车老厂房植入现代 Loft 钢结构与红砖）。",
    "  - 策略 3：公共空间激活（脏乱废弃边角地重塑为口袋公园）。"
])

# ==========================================
# 🚀 第六页：总结与反思 (P6)
# ==========================================
add_content_slide(prs, "5. 创新治理机制与课题总结", [
    "🧠 ‘规划师-公众-AI’协同模式：",
    "  - 引入大语言模型处理公众诉求，系统自动输出功能置换建议，实现全生命周期反馈管理。",
    "🎓 专业反思：",
    "  - AI 无法替代规划师的‘空间伦理判断’与人文关怀。",
    "  - 技术是赋能工具，最终的价值导向仍需依靠对‘古今共振’理念的深刻把握。"
])

# ==========================================
# 🚀 保存输出
# ==========================================
prs.save('毕设答辩_陈礼冲.pptx')
print("✅ 报告，PPT 骨架已成功生成！文件名为：毕设答辩_陈礼冲.pptx")