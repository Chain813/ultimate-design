# tools/generate_atlas_ppt.py
"""
Generate the atlas PPTX with:
  - Cover + TOC image slides
  - Chapter navigation divider slides (×6) before each chapter
  - All atlas image slides
  - A comprehensive thank-you / acknowledgments slide at the end
"""
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Mm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageFilter

# Force UTF-8 stdout
sys.stdout.reconfigure(encoding='utf-8')

ROOT = Path(__file__).resolve().parent.parent
ATLAS_DIR = ROOT / "static" / "atlas"
OUTPUT_DIR = ROOT / "static" / "附件"
OUTPUT_FILE = OUTPUT_DIR / "成果图册汇总_A3.pptx"

# Add scratch to path to import CHAPTERS
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "scratch"))

try:
    from gen_ppt import CHAPTERS
except ImportError:
    CHAPTERS = []

# Filter CHAPTERS to only contain existing files
_filtered = []
for ch_name, ch_en, sheets in CHAPTERS:
    _existing = []
    for fn, title in sheets:
        if (ATLAS_DIR / fn).exists():
            _existing.append((fn, title))
        else:
            print(f"  [OMIT FROM PPT OUTLINE] {fn}")
    _filtered.append((ch_name, ch_en, _existing))
CHAPTERS = _filtered


# ── Design tokens ──────────────────────────────────────────────
SLIDE_W = Mm(420)
SLIDE_H = Mm(297)

ORANGE = RGBColor(0xD4, 0x7C, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

CHAPTER_COLORS = [
    RGBColor(0xD4, 0x7C, 0x2A),  # Ch1 Orange
    RGBColor(0x3B, 0x82, 0xF6),  # Ch2 Blue
    RGBColor(0x8B, 0x5C, 0xF6),  # Ch3 Purple
    RGBColor(0x10, 0xB9, 0x81),  # Ch4 Green
    RGBColor(0xEF, 0x44, 0x44),  # Ch5 Red
    RGBColor(0x06, 0xB6, 0xD4),  # Ch6 Cyan
]

CHAPTER_EN = {
    "第1章 项目背景与概况": "PROJECT BACKGROUND",
    "第2章 现状调查与分析": "SITE INVESTIGATION",
    "第3章 设计理念与构思": "CONCEPT & STRATEGY",
    "第4章 总体方案设计":   "MASTER PLAN DESIGN",
    "第5章 重点地块设计":   "KEY PLOT DESIGN",
    "第6章 技术支撑":       "TECHNICAL SUPPORT",
}

# ── Helpers ─────────────────────────────────────────────────────

def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """Add a textbox with a single run of styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    
    # Enable East Asian / Complex Script fonts in XML for PowerPoint rendering
    from pptx.oxml.ns import qn
    from lxml import etree
    rPr = run._r.get_or_add_rPr()
    for tag in ['latin', 'ea', 'cs']:
        el = rPr.find(qn(f'a:{tag}'))
        if el is None:
            el = etree.SubElement(rPr, qn(f'a:{tag}'))
        el.set('typeface', font_name)
        
    return txBox


def _add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_centered_picture(slide, img_path, slide_width, slide_height):
    """Fit and center an image on the slide, maintaining aspect ratio."""
    with Image.open(img_path) as img:
        img_w, img_h = img.size

    img_aspect = img_w / img_h
    slide_aspect = slide_width / slide_height

    if img_aspect > slide_aspect:
        fit_width = slide_width
        fit_height = slide_width / img_aspect
    else:
        fit_height = slide_height
        fit_width = slide_height * img_aspect

    left = (slide_width - fit_width) / 2
    top = (slide_height - fit_height) / 2

    slide.shapes.add_picture(
        str(img_path),
        left=int(left), top=int(top),
        width=int(fit_width), height=int(fit_height)
    )


# ── Chapter Navigation Slide ───────────────────────────────────

def add_chapter_nav_slide(prs, blank_layout, active_chapter_idx):
    """
    Create a chapter navigation/divider slide.
    Left panel: TOC with all chapters (active one highlighted).
    Right panel: active chapter detail with sheet listing.
    """
    slide = prs.slides.add_slide(blank_layout)
    accent = CHAPTER_COLORS[active_chapter_idx % len(CHAPTER_COLORS)]

    # Left accent stripe (narrow)
    _add_rect(slide, 0, 0, Mm(4), SLIDE_H, accent)
    # Left panel background
    _add_rect(slide, Mm(4), 0, Mm(128), SLIDE_H, BG_LIGHT)

    # "目录" title
    _add_textbox(slide, Mm(14), Mm(12), Mm(110), Mm(16),
                 "目录", font_size=32, bold=True, color=BLACK)
    _add_textbox(slide, Mm(14), Mm(30), Mm(110), Mm(10),
                 "PRESENTATION OUTLINE", font_size=14, bold=False, color=MID_GRAY)

    # Divider line below title
    _add_rect(slide, Mm(14), Mm(42), Mm(104), Mm(0.8), accent)

    # Chapter list on left
    y_start = Mm(52)
    ch_height = Mm(36)

    for ci, (ch_name, ch_en_name, sheets) in enumerate(CHAPTERS):
        y_pos = y_start + ci * ch_height
        en_label = CHAPTER_EN.get(ch_name, ch_en_name)
        sheet_count = len(sheets)

        if ci == active_chapter_idx:
            # Highlight bar for active chapter
            _add_rect(slide, Mm(14), y_pos, Mm(3), Mm(24), accent)
            _add_rect(slide, Mm(20), y_pos - Mm(2), Mm(107), Mm(28),
                      RGBColor(0xFF, 0xFF, 0xFF))
            _add_textbox(slide, Mm(24), y_pos, Mm(100), Mm(14),
                         f">>> {ch_name}", font_size=16, bold=True, color=accent)
        else:
            _add_textbox(slide, Mm(24), y_pos, Mm(100), Mm(14),
                         ch_name, font_size=15, bold=False, color=DARK_TEXT)

        _add_textbox(slide, Mm(24), y_pos + Mm(15), Mm(100), Mm(10),
                     f"{en_label}  |  {sheet_count} sheets",
                     font_size=10, bold=False, color=MID_GRAY)

    # ── Right panel: active chapter detail ──
    ch_name, ch_en_name, sheets = CHAPTERS[active_chapter_idx]
    en_label = CHAPTER_EN.get(ch_name, ch_en_name)

    # Chapter title (large)
    _add_textbox(slide, Mm(148), Mm(12), Mm(248), Mm(20),
                 ch_name, font_size=36, bold=True, color=BLACK)
    _add_textbox(slide, Mm(148), Mm(36), Mm(248), Mm(12),
                 en_label, font_size=18, bold=False, color=MID_GRAY)

    # Accent line
    _add_rect(slide, Mm(148), Mm(54), Mm(248), Mm(0.8), accent)

    # Sheet count label
    _add_textbox(slide, Mm(148), Mm(60), Mm(248), Mm(10),
                 f"本章共 {len(sheets)} 张图纸:",
                 font_size=13, bold=False, color=DARK_TEXT)

    # Dynamically scale font size and column widths for chapters with many sheets
    if len(sheets) > 30:
        col_cap = 23
        font_size = 8
        dy = 9
        col_w = 62
    elif len(sheets) > 15:
        col_cap = 15
        font_size = 10
        dy = 13
        col_w = 85
    else:
        col_cap = 13
        font_size = 12
        dy = 16
        col_w = 125

    col_width = Mm(col_w)
    y_base = Mm(74)

    for si, (fn, title) in enumerate(sheets):
        col = si // col_cap
        row = si % col_cap
        x_pos = Mm(148) + col * col_width
        y_pos = y_base + row * Mm(dy)

        label = f"{si+1:02d}.  {title}"
        _add_textbox(slide, x_pos, y_pos, col_width, Mm(dy - 2),
                     label, font_size=font_size, bold=False, color=DARK_TEXT)

    print(f"  + Chapter nav slide: {ch_name}")
    return slide


# ── Thank You / Acknowledgments Slide ──────────────────────────

def add_thank_you_slide(prs, blank_layout):
    """
    Comprehensive acknowledgments slide with:
    - Advisor credits
    - AI model credits
    - Open source libraries & GitHub authors
    - Faculty acknowledgment
    """
    slide = prs.slides.add_slide(blank_layout)
    accent = RGBColor(0xD4, 0x7C, 0x2A)

    # Left accent stripe
    _add_rect(slide, 0, 0, Mm(4), SLIDE_H, accent)

    # Main title
    _add_textbox(slide, Mm(40), Mm(18), Mm(340), Mm(28),
                 "致  谢", font_size=48, bold=True, color=BLACK,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Mm(40), Mm(48), Mm(340), Mm(14),
                 "ACKNOWLEDGMENTS", font_size=20, bold=False, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Divider
    _add_rect(slide, Mm(170), Mm(66), Mm(80), Mm(1), accent)

    # ── Section 1: 指导教师 ──
    y = Mm(76)
    _add_textbox(slide, Mm(30), y, Mm(360), Mm(12),
                 "▎指导教师", font_size=20, bold=True, color=accent)
    _add_textbox(slide, Mm(40), y + Mm(16), Mm(340), Mm(10),
                 "崔诚慧 老师  ·  李冰心 老师",
                 font_size=16, bold=False, color=DARK_TEXT)
    _add_textbox(slide, Mm(40), y + Mm(28), Mm(340), Mm(10),
                 "感谢两位老师在选题、技术路线与设计方案全程中给予的悉心指导和宝贵意见。",
                 font_size=12, bold=False, color=MID_GRAY)

    # ── Section 2: AI 模型与工具 ──
    y = Mm(116)
    _add_textbox(slide, Mm(30), y, Mm(360), Mm(12),
                 "▎AI 模型与智能工具", font_size=20, bold=True, color=accent)
    ai_models = (
        "Google Gemini  ·  OpenAI ChatGPT  ·  Anthropic Claude  ·  DeepSeek  ·  Gamma  ·  "
        "NVIDIA SegFormer  ·  Stable Diffusion  ·  Ollama Gemma3"
    )
    _add_textbox(slide, Mm(40), y + Mm(16), Mm(340), Mm(14),
                 ai_models, font_size=13, bold=False, color=DARK_TEXT)
    _add_textbox(slide, Mm(40), y + Mm(32), Mm(340), Mm(10),
                 "以上模型在数据分析、语义分割、方案生成、文档编纂与合规审计等环节提供了关键技术支持。",
                 font_size=12, bold=False, color=MID_GRAY)

    # ── Section 3: 开源库与 GitHub 致谢 ──
    y = Mm(158)
    _add_textbox(slide, Mm(30), y, Mm(360), Mm(12),
                 "▎开源技术栈与 GitHub 社区", font_size=20, bold=True, color=accent)

    # Two columns of open source libs
    libs_col1 = [
        "Streamlit — Snowflake Inc.",
        "GeoPandas — Kelsey Jordahl et al.",
        "PyTorch — Meta AI (FAIR)",
        "Transformers — Hugging Face",
        "NumPy / SciPy — NumPy Community",
        "Pandas — Wes McKinney et al.",
        "Pillow (PIL) — Jeffrey A. Clark",
        "Plotly — Plotly Inc.",
        "Folium — Rob Story / Mapbox",
        "python-pptx — Steve Canny",
    ]
    libs_col2 = [
        "python-docx — Steve Canny",
        "Shapely — Sean Gillies et al.",
        "OpenStreetMap — OSM Contributors",
        "Jieba — Junyi Sun (fxsjy)",
        "BeautifulSoup4 — Leonard Richardson",
        "Selenium — SeleniumHQ",
        "PyMuPDF — Artifex Software",
        "pydeck — Uber / vis.gl",
        "Mammoth — Michael Williamson",
        "Ruff / Pre-commit — Astral / Yelp",
    ]

    for ri, lib in enumerate(libs_col1):
        _add_textbox(slide, Mm(40), y + Mm(16) + ri * Mm(10), Mm(160), Mm(9),
                     f"•  {lib}", font_size=11, bold=False, color=DARK_TEXT)

    for ri, lib in enumerate(libs_col2):
        _add_textbox(slide, Mm(210), y + Mm(16) + ri * Mm(10), Mm(170), Mm(9),
                     f"•  {lib}", font_size=11, bold=False, color=DARK_TEXT)

    # ── Section 4: 建筑与规划学院 ──
    y = Mm(268)
    _add_rect(slide, Mm(140), y - Mm(4), Mm(140), Mm(0.6), LIGHT_GRAY)
    _add_textbox(slide, Mm(40), y, Mm(340), Mm(10),
                 "感谢所有参与项目指导与支持的老师与专家",
                 font_size=14, bold=False, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Mm(40), y + Mm(12), Mm(340), Mm(10),
                 f"汇报人: {get_author_info().get('name','')}  |  {get_institution_info().get('name','')} {get_institution_info().get('department','')}".strip(),
                 font_size=12, bold=False, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

    print("  + Thank-you / Acknowledgments slide added.")
    return slide


# ── Main PPT Generation ────────────────────────────────────────

def generate_ppt():
    print("\nInitializing PPTX presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    slide_count = 0
    missing_files = []

    # ── 1) Cover + TOC ──
    for cover_fn in ["DR-001_规划设计图册封面.png", "DR-002_图册目录.png"]:
        img_path = ATLAS_DIR / cover_fn
        if img_path.exists():
            slide = prs.slides.add_slide(blank_layout)
            add_centered_picture(slide, img_path, prs.slide_width, prs.slide_height)
            slide_count += 1
            print(f"[{slide_count}] Cover/TOC: {cover_fn}")
        else:
            print(f"Warning: {cover_fn} not found. Skipping.")
            missing_files.append(cover_fn)

    # ── 2) Per-chapter: nav slide + image slides ──
    for ch_idx, (ch_name, ch_en_name, sheets) in enumerate(CHAPTERS):
        # Chapter navigation divider
        add_chapter_nav_slide(prs, blank_layout, ch_idx)
        slide_count += 1

        # Image slides for this chapter
        for fn, title in sheets:
            if fn in ["DR-001_规划设计图册封面.png", "DR-002_图册目录.png"]:
                continue
            img_path = ATLAS_DIR / fn
            if not img_path.exists():
                print(f"  Warning: [{fn}] not found. Skipping...")
                missing_files.append(fn)
                continue

            slide = prs.slides.add_slide(blank_layout)
            add_centered_picture(slide, img_path, prs.slide_width, prs.slide_height)
            slide_count += 1
            print(f"  [{slide_count}] {fn}")

    # ── 3) Thank you slide ──
    add_thank_you_slide(prs, blank_layout)
    slide_count += 1

    # ── Save to both destinations ──
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    paths_to_save = [
        OUTPUT_FILE,
        ATLAS_DIR / "答辩PPT.pptx"
    ]
    
    for out_path in paths_to_save:
        try:
            out_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(out_path))
            print(f"Success! PPTX generated at: {out_path}")
        except PermissionError:
            alt_output = str(out_path).replace(".pptx", "_new.pptx")
            prs.save(alt_output)
            print(f"[WARNING] Permission denied on {out_path}. Saved to: {alt_output}")
            
    print("\n" + "=" * 50)

    print(f"Total slides: {slide_count}")
    if missing_files:
        print(f"Missing files ({len(missing_files)}): {missing_files}")
    print("=" * 50)


if __name__ == "__main__":
    generate_ppt()
