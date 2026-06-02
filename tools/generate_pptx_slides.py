import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

# Separate template and output paths to prevent modifying template directly
user_home = os.path.expanduser("~")
ppt_template_path = os.path.join(user_home, "Desktop", "城环杯", "附件4 成果演示幻灯（模板）.pptx")
ppt_output_path = os.path.join(user_home, "Desktop", "城环杯", "附件4 成果演示幻灯.pptx")
temp_img_dir = os.path.join(user_home, "Desktop", "城环杯", "temp_images")
atlas_dir = r"e:\AI-based-project\urban-platform\static\atlas"
brain_prev_dir = os.path.join(user_home, ".gemini", "antigravity", "brain", "a7a0a585-8fe2-47a0-8b18-0be8b3147e91")

os.makedirs(temp_img_dir, exist_ok=True)

# Theme Colors (Modern Business/Planning Palette)
C_PRIMARY = RGBColor(24, 43, 73)        # Deep Navy Blue
C_ACCENT = RGBColor(0, 150, 136)        # Teal Green
C_TEXT_DARK = RGBColor(60, 64, 70)      # Charcoal Text
C_BG_CARD = RGBColor(250, 252, 255)     # Off-White Card Fill
C_BORDER_CARD = RGBColor(220, 226, 235) # Card border
C_RED_WARN = RGBColor(220, 53, 69)       # Crimson Red (Warning)
C_GREEN_OK = RGBColor(40, 167, 69)       # Emerald Green (Pass)

print("Loading PPTX template...")
prs = Presentation(ppt_template_path)

def format_title_and_add_bar(slide, text):
    # Format and reposition title shape
    title_shape = slide.shapes.title
    title_shape.left = Inches(1.1)
    title_shape.top = Inches(0.4)
    title_shape.width = Inches(11.0)
    title_shape.height = Inches(0.6)
    title_shape.text = text
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = "微软雅黑"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.alignment = PP_ALIGN.LEFT
    title_shape.text_frame.margin_left = Inches(0)
    title_shape.text_frame.margin_top = Inches(0)
    title_shape.text_frame.margin_bottom = Inches(0)
    title_shape.text_frame.margin_right = Inches(0)
    
    # Add a thin vertical teal bar on the left of the title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.48), Inches(0.08), Inches(0.36))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

def create_card_with_left_bar(slide, left, top, width, height, fill_color=C_BG_CARD, border_color=C_BORDER_CARD, bar_color=C_ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    
    # Draw left overlay vertical bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.02), top + Inches(0.06), Inches(0.06), height - Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    return card

def create_tabbed_card(slide, left, top, width, height, fill_color=C_BG_CARD, border_color=C_BORDER_CARD, bar_color=C_ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    
    # Draw top overlay horizontal bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.06), top + Inches(0.02), width - Inches(0.12), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    return card

def add_card_text_runs(card, title_text, bullets_data):
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.15)
    
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "微软雅黑"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = C_PRIMARY
    p_title.space_after = Pt(6)
    
    for bold_part, regular_part in bullets_data:
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = 1.15
        
        r_bold = p.add_run()
        r_bold.text = bold_part
        r_bold.font.name = "微软雅黑"
        r_bold.font.size = Pt(10.5)
        r_bold.font.bold = True
        r_bold.font.color.rgb = C_ACCENT
        
        r_reg = p.add_run()
        r_reg.text = regular_part
        r_reg.font.name = "微软雅黑"
        r_reg.font.size = Pt(10.5)
        r_reg.font.color.rgb = C_TEXT_DARK

def add_kpi_card_new(slide, left, top, width, height, value_text, label_text, is_warn=False):
    # Card background (white)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(220, 225, 235)
    card.line.width = Pt(1)
    
    # Left vertical accent line
    accent_color = C_RED_WARN if is_warn else C_GREEN_OK
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    # Text frame
    txBox = slide.shapes.add_textbox(left + Inches(0.12), top, width - Inches(0.12), height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)
    
    p_val = tf.paragraphs[0]
    p_val.text = value_text
    p_val.font.name = "Arial"
    p_val.font.size = Pt(28)
    p_val.font.bold = True
    p_val.font.color.rgb = accent_color
    p_val.alignment = PP_ALIGN.LEFT
    p_val.space_after = Pt(2)
    
    p_lbl = tf.add_paragraph()
    p_lbl.text = label_text
    p_lbl.font.name = "微软雅黑"
    p_lbl.font.size = Pt(10)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = C_TEXT_DARK
    p_lbl.alignment = PP_ALIGN.LEFT

def add_framed_picture(slide, img_path, left, top, width, height):
    from PIL import Image
    with Image.open(img_path) as img:
        img_w, img_h = img.size
    
    img_aspect = img_w / img_h
    box_aspect = width / height
    
    if img_aspect > box_aspect:
        # Image is wider than container -> constrain by width
        scaled_width = width
        scaled_height = width / img_aspect
    else:
        # Image is taller than container -> constrain by height
        scaled_height = height
        scaled_width = height * img_aspect
        
    offset_x = (width - scaled_width) / 2
    offset_y = (height - scaled_height) / 2
    
    frame_left = left + offset_x
    frame_top = top + offset_y
    
    # White background frame shape fitting the scaled image
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_left, frame_top, int(scaled_width), int(scaled_height))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
    frame.line.color.rgb = RGBColor(210, 215, 225)
    frame.line.width = Pt(1)
    
    # Insert picture inside frame with a slight inset margin
    inset = Inches(0.04)
    if scaled_width - inset*2 > 0 and scaled_height - inset*2 > 0:
        slide.shapes.add_picture(
            img_path, 
            frame_left + inset, 
            frame_top + inset, 
            width=int(scaled_width - inset*2), 
            height=int(scaled_height - inset*2)
        )
    else:
        slide.shapes.add_picture(
            img_path, 
            frame_left, 
            frame_top, 
            width=int(scaled_width), 
            height=int(scaled_height)
        )

# Update Slide 1 (Cover)
print("Updating Slide 1 Cover...")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame and "第十届" in shape.text_frame.text:
        shape.text_frame.text = "城垣杯\n规划决策支持模型设计大赛"
        for p in shape.text_frame.paragraphs:
            p.font.name = "黑体"
            p.font.size = Pt(40)
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
    elif shape.has_text_frame and "The 10th" in shape.text_frame.text:
        shape.text_frame.text = "The Planning Decision Support Model Design Contest (Chengyuan Cup)"
        for p in shape.text_frame.paragraphs:
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT

# Update Slide 2 (Registration Info)
print("Updating Slide 2 Info...")
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "研究题目" in txt:
            shape.text_frame.text = "基于大模型与多模态AI的城市更新\n空间设计智能推演系统"
            shape.width = Inches(10.0)
            shape.height = Inches(2.2)
            shape.top = Inches(0.25)
            # Resize the orange background block (Shape 3) to stretch it horizontally
            if len(slide2.shapes) > 3:
                s3 = slide2.shapes[3]
                s3.width = Inches(10.5)
                s3.height = Inches(2.5)
            for p in shape.text_frame.paragraphs:
                p.font.name = "微软雅黑"
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.LEFT

        elif "参赛者：" in txt:
            shape.text_frame.text = "参赛者：陈礼冲、刘旭东\n参赛单位：吉林建筑大学\n投稿方向：主题二：面向高质量发展的城市治理\n报名编号：F221"
            for p in shape.text_frame.paragraphs:
                p.font.name = "微软雅黑"
                p.font.size = Pt(18)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
        elif "第十届" in txt:
            shape.text_frame.text = "城垣杯 • 规划决策支持模型设计大赛"
            for p in shape.text_frame.paragraphs:
                p.font.name = "黑体"
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.LEFT

# Update Slide 10 (Ending slide)
print("Updating Slide 10 Ending...")
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if shape.has_text_frame and "本模板仅供参考" in shape.text_frame.text:
        shape.text_frame.text = "谢谢大家！\n请各位专家批评指正"
        p = shape.text_frame.paragraphs[0]
        p.font.name = "黑体"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Create layouts using Title Only Layout (Layout 5)
layout_title_only = prs.slide_layouts[5]

# --- NEW SLIDE 1 (Background) ---
print("Creating Slide: Background...")
s1 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s1, "研究背景：存量更新时代城市治理三维挑战")

# Left: 3 vertical cards with warning left bar
card_w = Inches(5.8)
card_h = Inches(1.5)
left_x = Inches(0.8)
gap_y = Inches(0.18)

c1 = create_card_with_left_bar(s1, left_x, Inches(1.4), card_w, card_h, bar_color=C_RED_WARN)
add_card_text_runs(c1, "01 / 城市体检诊断粗糙 (Diagnosis Gap)", [
    ("细粒度不足：", "传统体检评估严重依赖宏观数据和专家主观定性分析，缺乏地块级、多源数据融合的定量体检方法。")
])

c2 = create_card_with_left_bar(s1, left_x, Inches(1.4) + card_h + gap_y, card_w, card_h, bar_color=C_RED_WARN)
add_card_text_runs(c2, "02 / 利益主体博弈断裂 (Consensus Gap)", [
    ("诉求冲突大：", "微更新涉及政府（控规红线）、开发商（商业回报）与居民（人居环境）的复杂博弈，共识达成困难。")
])

c3 = create_card_with_left_bar(s1, left_x, Inches(1.4) + (card_h + gap_y)*2, card_w, card_h, bar_color=C_RED_WARN)
add_card_text_runs(c3, "03 / AIGC 制图幻觉严重 (Design Gap)", [
    ("空间失位：", "普通生成式 AI 直接用于规划制图时容易产生地理空间拓扑错乱和“幻觉”，无法对齐法定红线指标。")
])

# Right: Framed map
img_path = os.path.join(temp_img_dir, "fig_017.png")
if os.path.exists(img_path):
    add_framed_picture(s1, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))


# --- NEW SLIDE 2 (Goal) ---
print("Creating Slide: Goal...")
s2 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s2, "研究目标：大模型与多模态AI驱动的城市更新平台")

c1_s2 = create_card_with_left_bar(s2, left_x, Inches(1.4), card_w, card_h, bar_color=C_ACCENT)
add_card_text_runs(c1_s2, "01 / AHP-MPI 空间体检诊断引擎", [
    ("定量诊断：", "融合潜力、需求与现状构建 MPI 指数，对研究区 719 栋现状建筑进行精细化、地块级更新潜力评估。")
])

c2_s2 = create_card_with_left_bar(s2, left_x, Inches(1.4) + card_h + gap_y, card_w, card_h, bar_color=C_ACCENT)
add_card_text_runs(c2_s2, "02 / 多智能体博弈决策协商沙盘", [
    ("共识达成：", "基于 LLM 多智能体设计居民、开发商、政府规划局角色，开发满意度效用算法，动态判定并达成共识。")
])

c3_s2 = create_card_with_left_bar(s2, left_x, Inches(1.4) + (card_h + gap_y)*2, card_w, card_h, bar_color=C_ACCENT)
add_card_text_runs(c3_s2, "03 / 空间对齐 AI 规划绘图管线", [
    ("图纸编译：", "设计「矢量-光栅-ControlNet」约束制图，一键并行编译输出 26 张 A3 标准规划图，消除空间位置幻觉。")
])

img_path = os.path.join(temp_img_dir, "system_architecture.png")
if os.path.exists(img_path):
    add_framed_picture(s2, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))



# --- NEW SLIDE 3 (Methodology) ---
print("Creating Slide: Methodology...")
s3 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s3, "核心研究方法与多模态数字理论框架")

# 3 Horizontal columns
col_w = Inches(3.7)
col_h = Inches(4.86)
col_y = Inches(1.4)
gap_x = Inches(0.3)

mc1 = create_tabbed_card(s3, Inches(0.8), col_y, col_w, col_h, bar_color=C_ACCENT)
add_card_text_runs(mc1, "AHP-MPI 空间体检模型", [
    ("因子权重测度：", "融合空间潜力(S)、配套需求(D)、环境质量(E)三维度，采用层次分析法(AHP)确定决策矩阵因子权重，一致性比率 CR < 0.1。"),
    ("潜力科学识别：", "MPI = (0.4S + 0.3D + 0.3(1.0-E))*100，利用 (1.0 - E) 代表环境越差，绿化织补需求度越高。")
])

mc2 = create_tabbed_card(s3, Inches(0.8) + col_w + gap_x, col_y, col_w, col_h, bar_color=C_PRIMARY)
add_card_text_runs(mc2, "多智能体博弈决策模型", [
    ("角色诉求扮演：", "构建居民（配套/绿化）、开发商（商业运营/容积率）、规划师（历史文化保护/限高）多智能体，实现智能谈判。"),
    ("满意度效用监控：", "S_role = min(100, 50 + 7*count(DialogueText 命中关键词))。当三方底线满意度均达到 60 分共识线时自动收敛。")
])

mc3 = create_tabbed_card(s3, Inches(0.8) + (col_w + gap_x)*2, col_y, col_w, col_h, bar_color=C_ACCENT)
add_card_text_runs(mc3, "空间句法与视觉绿化分析", [
    ("路网拓扑分析：", "基于开源 OSMnx 与 NetworkX 库对 74 段核心路段建模，计算全局整合度与穿行度，识别交通连通织补瓶颈。"),
    ("视觉绿视率：", "应用 SegFormer 开源模型批量分割 1,788 张街景影像，解析各样点平均绿视率(GVI)，作为环境品质输入。")
])


# --- NEW SLIDE 4 (Route/Pipeline) ---
print("Creating Slide: Pipeline...")
s4 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s4, "系统技术路线与 GIS-to-AIGC 空间对齐管线")

c1_s4 = create_card_with_left_bar(s4, left_x, Inches(1.4), card_w, Inches(2.35), bar_color=C_PRIMARY)
add_card_text_runs(c1_s4, "1. 闭环决策技术路线", [
    ("四引擎闭环工作流：", "打通“多源数据体检诊断 -> 三主体博弈协商 -> RAG国家法规合规审查 -> AIGC自动图册编译”，构成紧凑的数字化决策环。"),
    ("前后端架构：", "Streamlit 1.55 进行全交互Web发布，算法引擎与GIS底层参数彻底分离，一键迁移到其它更新片区。")
])

c2_s4 = create_card_with_left_bar(s4, left_x, Inches(1.4) + Inches(2.35) + gap_y, card_w, Inches(2.35), bar_color=C_ACCENT)
add_card_text_runs(c2_s4, "2. AIGC 空间对齐与 Stable Diffusion", [
    ("矢量转光栅控制：", "将用地和路网 GeoJSON 转换为分类色彩光栅底图，作为 ControlNet 的输入硬约束。"),
    ("消解地理空间幻觉：", "融合 Canny (边缘骨架) 与 Seg (用地红线) 双重 ControlNet，在 Diffusers 开源底座上确保图纸实现像素级对齐。")
])

img_path = os.path.join(temp_img_dir, "negotiation_workflow.png")
if os.path.exists(img_path):
    add_framed_picture(s4, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))



# --- NEW SLIDE 5 (Data) ---
print("Creating Slide: Data...")
s5 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s5, "数据说明：研究区多源空间数据库构建")

# Left Column: Large metric dashboard cards (High contrast)
kpi1 = create_card_with_left_bar(s5, Inches(0.8), Inches(1.4), Inches(2.8), Inches(1.1), bar_color=C_PRIMARY)
add_card_text_runs(kpi1, "研究区范围 (Boundary)", [("170.2 公顷", " (自绘边界)")])

kpi2 = create_card_with_left_bar(s5, Inches(3.8), Inches(1.4), Inches(2.8), Inches(1.1), bar_color=C_ACCENT)
add_card_text_runs(kpi2, "现状建筑数量 (Buildings)", [("719 栋", " (OSM/实地核对)")])

kpi3 = create_card_with_left_bar(s5, Inches(0.8), Inches(2.65), Inches(2.8), Inches(1.1), bar_color=C_ACCENT)
add_card_text_runs(kpi3, "拓扑分析路网 (Roads)", [("74 段", " (OSM路网切分)")])

kpi4 = create_card_with_left_bar(s5, Inches(3.8), Inches(2.65), Inches(2.8), Inches(1.1), bar_color=C_PRIMARY)
add_card_text_runs(kpi4, "街景全景图像 (Panoramas)", [("1,788 张", " (百度街景采样)")])

c_s5_desc = create_card_with_left_bar(s5, Inches(0.8), Inches(3.9), Inches(5.8), Inches(2.36), bar_color=C_PRIMARY)
add_card_text_runs(c_s5_desc, "感知配套与法规数据库构建", [
    ("多源感知设施：", "抓取 411 条 POI 设施点计算核密度服务；收集 207 条微博舆情数据进行公众更新意向情感文本分析。"),
    ("国家保护规章向量化：", "切分 7 份保护规划法规及地方条例为 248 个高维语义块，导入轻量本地向量库，支持 RAG 合规实时校验。")
])

# Right Column: Styled Table
tb_left = Inches(6.9)
tb_top = Inches(1.4)
tb_width = Inches(5.6)
tb_height = Inches(4.86)

table_shape = s5.shapes.add_table(6, 4, tb_left, tb_top, tb_width, tb_height)
table = table_shape.table
t_data = [
    ["数据类别", "数据格式", "数据规模", "模型应用作用"],
    ["研究范围", "GeoJSON", "170.2 公顷", "确立规划边界"],
    ["现状建筑", "GeoJSON", "719 栋", "提取层数计算FAR"],
    ["现状路网", "GeoJSON", "74 段", "空间句法拓扑分析"],
    ["实景图像", "JPG", "1,788 张", "街景GVI语义分割"],
    ["政策法规", "Vector", "248分块", "RAG合规校核"]
]

for r_idx, row in enumerate(t_data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.name = "微软雅黑"
        p.font.size = Pt(10)
        p.font.bold = (r_idx == 0)
        p.alignment = PP_ALIGN.CENTER
        if r_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_PRIMARY
            p.font.color.rgb = RGBColor(255, 255, 255)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(252, 253, 254) if r_idx % 2 == 0 else RGBColor(245, 247, 250)


# --- NEW SLIDE 6 (Prep) ---
print("Creating Slide: Prep...")
s6 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s6, "数据预处理：投影纠偏与空间属性挂接")

c1_s6 = create_card_with_left_bar(s6, left_x, Inches(1.4), card_w, Inches(2.35), bar_color=C_PRIMARY)
add_card_text_runs(c1_s6, "1. 本地 Gauss-Kruger 投影纠偏 (EPSG:32651)", [
    ("高纬度拉伸形变消除：", "全国或全球性坐标系在长春高纬度地区存在极高形变差。统一重投影为长春本地 Gauss-Kruger EPSG:32651 三度带坐标。"),
    ("面积精确锁定：", "消除了 Web 墨卡托在中高纬度高达 93% 的面积计算拉伸形变，将研究边界精准锁定为 170.2公顷。")
])

c2_s6 = create_card_with_left_bar(s6, left_x, Inches(1.4) + Inches(2.35) + gap_y, card_w, Inches(2.35), bar_color=C_ACCENT)
add_card_text_runs(c2_s6, "2. 基于 SegFormer 的绿视率 (GVI) 测度", [
    ("像素级语义分割：", "SegFormer (ADE20K) 开源神经网络对 447 个物理样点的 1,788 张街景全景图进行四方向植被像素自动分割提取。"),
    ("识别生态短板：", "计算出全域平均绿视率 (GVI) 仅为 8.7%，从视觉感知角度坐实绿化极度匮乏事实，输入 MPI 数据集。")
])

img_path = os.path.join(temp_img_dir, "fig_plotly.png")
if os.path.exists(img_path):
    add_framed_picture(s6, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))


# --- NEW SLIDE 7 (Alg 1) ---
print("Creating Slide: Alg 1...")
s7 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s7, "算法原理：潜力指数与三主体满意度效用函数")

# Two wide horizontal columns
col_w_wide = Inches(5.6)
col_h_wide = Inches(4.86)

c1_s7 = create_tabbed_card(s7, Inches(0.8), Inches(1.4), col_w_wide, col_h_wide, bar_color=C_ACCENT)
add_card_text_runs(c1_s7, "AHP-MPI 空间体检模型", [
    ("MPI 潜力指数公式：", "MPI_i = (0.4 * S_i + 0.3 * D_i + 0.3 * (1.0 - E_i)) * 100\n"),
    ("变量解析：", "S_i 代表空间潜力（现状层数倒数与基底面积归一化）；D_i 代表社会需求（设施配套 POI 在 150m 半径内求核密度）；E_i 代表环境品质（街景 GVI 分数值）。"),
    ("公式亮点：", "采用 (1.0 - E_i) 表征当前绿化越差，其微更新的迫切性及绿化织补需求度越高。")
])

c2_s7 = create_tabbed_card(s7, Inches(6.8), Inches(1.4), col_w_wide, col_h_wide, bar_color=C_PRIMARY)
add_card_text_runs(c2_s7, "三主体满意度效用与共识判定", [
    ("满意度效用函数：", "S_role = min(100, 50 + 7 * count(DialogueText 命中关键词))\n"),
    ("关键词监控：", "居民 K_res = [“绿”, “公园”, “配套”, “菜市”, “口袋”]；开发商 K_dev = [“容积率”, “收益”, “文旅”, “回报”]；规划局 K_gov = [“历史保护”, “紫线”, “限高”, “条例”]。"),
    ("共识收敛机制：", "设定共识底线 min(S_res, S_dev, S_gov) >= 60。若任意一方满意度低于 60 分则触发黄色警报，大模型自动调整谈判策略，引导妥协直至共识达成。")
])


# --- NEW SLIDE 8 (Alg 2) ---
print("Creating Slide: Alg 2...")
s8 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s8, "合规校验：空间控规审查与系统架构")

c1_s8 = create_card_with_left_bar(s8, left_x, Inches(1.4), card_w, Inches(2.35), bar_color=C_PRIMARY)
add_card_text_runs(c1_s8, "1. Zoning Compliance 控规实时校验", [
    ("指标平面演算：", "FAR = sum( Floor_i * Area_i ) / Area_land, Density = sum( Area_footprint ) / Area_land。"),
    ("超标红牌告警：", "系统将空间计算数值与控规标准限值进行拓扑包含校验（容积率≤1.40，最高限高≤18.0m，核心区≤9.0m）。超出上限自动高亮警告。")
])

c2_s8 = create_card_with_left_bar(s8, left_x, Inches(1.4) + Inches(2.35) + gap_y, card_w, Inches(2.35), bar_color=C_ACCENT)
add_card_text_runs(c2_s8, "2. GitHub 开源核心技术栈", [
    ("空间几何分析：", "基于 GeoPandas 1.0 和 Shapely 2.0 开源库处理 719 栋建筑和 108 宗地块的几何拓扑校验；"),
    ("路网拓扑分析：", "调用 OSMnx 1.9 与 NetworkX 3.2 对 74 段核心路段进行连通性拓扑分析；"),
    ("大模型与 AI 绘图：", "后台部署 DeepSeek-V4 大模型进行博弈推理，前端交互采用 Streamlit 1.55。")
])

img_path = os.path.join(temp_img_dir, "compliance_audit_flow.png")
if os.path.exists(img_path):
    add_framed_picture(s8, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))



# --- NEW SLIDE 9 (Case 1) ---
print("Creating Slide: Case 1...")
s9 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s9, "应用实证：长春协调更新片区现状体检诊断")

# Left Column: Detailed description
c1_s9 = create_card_with_left_bar(s9, left_x, Inches(1.4), card_w, Inches(4.86), bar_color=C_PRIMARY)
add_card_text_runs(c1_s9, "历史与现状环境体检报告", [
    ("地块及建筑底数：", "长春宽城区协调片区占地 170.2 公顷，共计 108 宗地块，现状建筑 719 栋，层数多在 1-4 层（均值 3.7 层）。"),
    ("容积率与建筑密度：", "实测容积率 1.13，建筑密度 30.0%。两项均合规通过（FAR≤1.40，密度≤35%），开发强度整体可控。"),
    ("绿地率与天际线视廊：", "实测现状绿地率仅 2.9%！远低于 25% 国家标准，生态赤字巨大。最高建筑达 59.5 米，超出 18m 限高，破坏历史街区整体风貌天际线。")
])

# Right Column: Dashboard KPI grids (Using modern left-accent widgets)
k_w = Inches(2.6)
k_h = Inches(2.25)
x1 = Inches(6.9)
x2 = Inches(9.8)
y1 = Inches(1.4)
y2 = Inches(4.0)

add_kpi_card_new(s9, x1, y1, k_w, k_h, "1.13", "现状容积率\n(控规≤1.40)\n✅ 合规达标", is_warn=False)
add_kpi_card_new(s9, x2, y1, k_w, k_h, "30.0%", "现状建筑密度\n(控规≤35.0%)\n✅ 合规达标", is_warn=False)
add_kpi_card_new(s9, x1, y2, k_w, k_h, "2.9%", "现状绿地率\n(控规≥25.0%)\n❌ 严重违规 (偏低)", is_warn=True)
add_kpi_card_new(s9, x2, y2, k_w, k_h, "59.5m", "最高建筑高度\n(控规≤18.0m)\n⚠️ 局部溢出 (超高)", is_warn=True)


# --- NEW SLIDE 10 (Case 2) ---
print("Creating Slide: Case 2...")
s10 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s10, "博弈协商：多主体利益协商与共识收敛")

c1_s10 = create_card_with_left_bar(s10, left_x, Inches(1.4), card_w, Inches(2.35), bar_color=C_PRIMARY)
add_card_text_runs(c1_s10, "1. 第一轮辩论冲突：黄色预警触发", [
    ("诉求分歧点：", "开发商追求最大利益，强推 36m 中高层开发，获得满意度 80 分。但这引起居民强烈抵触（日照/压迫，满意度 45 分），规划局亦因风貌违背打出 65 分。"),
    ("预警触发：", "min(满意度) = 45 < 60，系统触发黄色利益冲突警报，并提出整改意见。")
])

c2_s10 = create_card_with_left_bar(s10, left_x, Inches(1.4) + Inches(2.35) + gap_y, card_w, Inches(2.35), bar_color=C_ACCENT)
add_card_text_runs(c2_s10, "2. 第二轮方案优化：共识协商达成", [
    ("妥协设计：", "规划局引导开发商，坚守核心区 18m 限高红线，在站前区做适度商业增容，并见缝插针地配套 3 处口袋公园以补偿居民绿化诉求。"),
    ("共识达成：", "优化后，三方满意度得分最终为：居民 66、开发商 73、规划局 68，均跨过 60 分共识底线。共识生成，自动提取设计导则。")
])

img_path = os.path.join(temp_img_dir, "fig_radar.png")
if os.path.exists(img_path):
    add_framed_picture(s10, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))


# --- NEW SLIDE 11 (Case 3) ---
print("Creating Slide: Case 3...")
s11 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s11, "成果表达：代码自动编译排版与高清A3图册")

c1_s11 = create_card_with_left_bar(s11, left_x, Inches(1.4), card_w, Inches(2.35), bar_color=C_PRIMARY)
add_card_text_runs(c1_s11, "1. 高清 A3 图框拼装排版流水线", [
    ("图框自动封装：", "系统利用 Matplotlib 绘制空间分析图层，配合 Pillow 对 A3 标准图框、图例框、线划比例尺及图签进行像素级拼接合并。"),
    ("高清批量导出：", "启用 Python 多进程并行模块，全自动渲染并导出包含现状图、控制性详细规划图、绿地与交通图在内的 26 张高清规划大图。")
])

c2_s11 = create_card_with_left_bar(s11, left_x, Inches(1.4) + Inches(2.35) + gap_y, card_w, Inches(2.35), bar_color=C_ACCENT)
add_card_text_runs(c2_s11, "2. 大模型图纸设计说明动态合成", [
    ("文图指标物理绑定：", "大图图签右下侧的设计说明是由大模型在拼版瞬间读取 GIS 物理数据库实测指标自动合成的，不存在手工录入延迟。"),
    ("消除人为不一致：", "保证了“规划图纸表现-数据库量化指标-设计文字说明”三者保持绝对逻辑一致。")
])

img_path = os.path.join(temp_img_dir, "fig_004.png")
if os.path.exists(img_path):
    add_framed_picture(s11, img_path, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.86))


# --- NEW SLIDE 12 (Summary) ---
print("Creating Slide: Summary...")
s12 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s12, "研究总结：平台核心技术创新与应用前景")

mc1_s12 = create_tabbed_card(s12, Inches(0.8), col_y, col_w, col_h, bar_color=C_PRIMARY)
add_card_text_runs(mc1_s12, "规划决策决策流程创新", [
    ("循证设计闭环：", "填补国内数字辅助微更新决策系统空白，首创集“体检诊断-多方博弈-合规校验-自动出图”于一体的端到端系统。"),
    ("民主决策底座：", "打破传统定性经验式判断，通过量化潜力测度与智能体谈判沙盘极大提高了决策的透明度与民主化水平。")
])

mc2_s12 = create_tabbed_card(s12, Inches(0.8) + col_w + gap_x, col_y, col_w, col_h, bar_color=C_ACCENT)
add_card_text_runs(mc2_s12, "AI 空间绘图技术创新", [
    ("消解空间位置幻觉：", "设计「矢量-光栅-ControlNet」双重控制通道，使得大模型生成的蓝图规划图能像素级精准贴合真实的规划红线边界。"),
    ("图文指标严密绑定：", "读取空间计算结果实时生成设计导则，摆脱传统手工编写的误差，多进程出图效率极高。")
])

mc3_s12 = create_tabbed_card(s12, Inches(0.8) + (col_w + gap_x)*2, col_y, col_w, col_h, bar_color=C_PRIMARY)
add_card_text_runs(mc3_s12, "工程应用与推广前景", [
    ("常态化体检：", "可直接对接到住建部门数据库，作为常态化体检工具，动态捕捉街区“绿化赤字”、“建筑违规”等事实。"),
    ("公众参与沙盘：", "作为街道办事处、居民与开发商的共识谈判平台，实时出图，高效沟通。"),
    ("教学辅助系统：", "可供城乡规划专业教学示范使用。")
])


# --- NEW SLIDE 13 (References & Open Source Stack) ---
print("Creating Slide: References & Open Source Stack...")
s13 = prs.slides.add_slide(layout_title_only)
format_title_and_add_bar(s13, "参考文献与 GitHub 开源技术底座")

# 2 Wide vertical cards
c1_s13 = create_tabbed_card(s13, Inches(0.8), Inches(1.4), col_w_wide, col_h_wide, bar_color=C_PRIMARY)
add_card_text_runs(c1_s13, "学术参考文献 (Academic References)", [
    ("[1] 丁梦月. ", "基于计算机视觉技术的城市街道步行空间人群行为原型研究[D].南京:东南大学,2021."),
    ("[2] 尧馨雅. ", "基于可解释深度学习的街道风貌基因图谱识别研究[D].杭州:浙江大学,2022."),
    ("[3] 赵卉. ", "历史肌理延续的数字化城市设计方法研究——以江苏同里古镇为例[D].南京:东南大学,2021."),
    ("[4] 卢文正. ", "社会空间理论视角的社区更新[D].哈尔滨:哈尔滨工业大学,2020."),
    ("[5] 张峰. ", "智慧城市空间信息资源规划的模型和实现方法研究[D].武汉:武汉大学,2005."),
    ("[6] 张国政. ", "数字孪生技术提升城市韧性路径研究[D].上海:华东政法大学,2023."),
    ("[7] 梁汉雄. ", "基于街景图片与深度学习的旧工业园区改造与可步行性要素研究[D].广州:华南理工大学,2022."),
    ("[8] 方可. ", "探索北京旧城居住区有机更新的适宜途径[D].北京:清华大学,2000.")
])

c2_s13 = create_tabbed_card(s13, Inches(6.8), Inches(1.4), col_w_wide, col_h_wide, bar_color=C_ACCENT)
add_card_text_runs(c2_s13, "GitHub 开源项目与技术栈致谢", [
    ("gboeing/osmnx (v1.9.0)：", "提供街区路网拓扑分析与空间句法整合度指标计算。"),
    ("huggingface/transformers (v4.38)：", "搭载预训练 SegFormer 语义分割大模型，提取环境感知指标。"),
    ("lllyasviel/ControlNet (Diffusers)：", "实现 Canny-Seg 空间物理双通道位置对齐控规红线制图管线。"),
    ("geopandas/geopandas (v1.0.1)：", "提供高效的 GeoJSON/Shapefile 空间矢量几何拓扑关系校核与投影。"),
    ("streamlit/streamlit (v1.55.0)：", "构建纯 Python 编写的 Web 交互决策沙盘与合规告警可视化底座。"),
    ("plotly/plotly.py (v5.24.0)：", "实现 WebGL 交互式 3D 建筑高度及博弈满意度雷达图渲染。")
])


# Reorder slides dynamically to place content immediately after respective transitions
print("Reordering slides...")
slide_id_list = prs.slides._sldIdLst
slide_ids = list(slide_id_list)

# The new order of slide XML element indices:
# Original template: 0:Cover, 1:Info, 2:TOC, 3:Transition 1, 4:Transition 2, 5:Transition 3, 6:Transition 4, 7:Transition 5, 8:Transition 6, 9:Ending
# Content slides (indices 10 to 22):
# 10:Background, 11:Goal, 12:Method, 13:Pipeline, 14:Data, 15:Prep, 16:Alg 1, 17:Alg 2, 18:Case 1, 19:Case 2, 20:Case 3, 21:Summary, 22:References
new_order_indices = [
    0, 1, 2,        # Cover, Info, TOC
    3, 10, 11,      # 一、研究问题 (Transition 1) -> Background, Goal
    4, 12, 13,      # 二、研究方法 (Transition 2) -> Method, Pipeline
    5, 14, 15,      # 三、数据说明 (Transition 3) -> Data, Prep
    6, 16, 17,      # 四、模型算法 (Transition 4) -> Alg 1, Alg 2
    7, 18, 19, 20,  # 五、实践案例 (Transition 5) -> Case 1, Case 2, Case 3
    8, 21, 22,      # 六、研究总结 (Transition 6) -> Summary, References
    9               # Ending slide
]

slide_id_list.clear()
for idx in new_order_indices:
    slide_id_list.append(slide_ids[idx])

# Save Presentation to final output, keeping template file clean, handling PermissionError
try:
    prs.save(ppt_output_path)
    print("PPTX presentation generated and saved successfully!")
except PermissionError:
    alt_path = ppt_output_path.replace(".pptx", "_更新.pptx")
    prs.save(alt_path)
    print(f"Permission denied on {ppt_output_path} (probably open in PowerPoint). Saved instead to: {alt_path}")
